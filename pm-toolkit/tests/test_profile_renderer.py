import pytest
import shutil
from pathlib import Path
from unittest.mock import patch
from core.profile_schema import ProfileSchema, ExperienceEntry, ProjectEntry

TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "base_profile_template.pptx"
SPEC_PATH = Path(__file__).parent.parent / "templates" / "template_spec.json"


FIXTURE_PROFILE = ProfileSchema(
    name="Virgil Galan",
    role_title="Technical Director | Agile Delivery | Cloud & Digital Transformation",
    role_subtitle="Technical Director",
    profile="Virgil Galan is a technical director with 15 years of experience leading digital transformation programmes across financial services and retail. Delivered multiple cloud migration and agile delivery programmes.",
    competencies=[
        "Agile Delivery", "Cloud Architecture", "Stakeholder Management",
        "Budget Control", "Risk Management", "Team Leadership",
        "Digital Transformation", "DevOps"
    ],
    education=["MSc Computer Science, Polytechnic University, 2005"],
    certifications=["PMP", "AWS Solutions Architect"],
    methodologies=["Scrum", "Kanban", "SAFe", "PRINCE2"],
    technologies=["AWS", "Azure", "Kubernetes", "Terraform", "Jira", "Confluence"],
    experience=[
        ExperienceEntry(
            employer="Cognizant Softvision",
            role="Technical Director",
            date_range="Jan 2018 – Present",
            employer_bullets=["Led a 60-person delivery organisation across 3 geographies."],
            projects=[
                ProjectEntry(
                    project_name="Hospitality – Hotel Booking Platform",
                    date_range="2022–2025",
                    bullets=[
                        "Architected microservices migration from monolithic legacy platform.",
                        "Reduced deployment lead time by 60% through CI/CD automation."
                    ]
                )
            ]
        )
    ]
)


@pytest.mark.skipif(
    not TEMPLATE_PATH.exists(),
    reason="base_profile_template.pptx not yet created (Phase 1 manual step)"
)
def test_renderer_produces_valid_pptx(tmp_path):
    from core.profile_renderer import ProfileRenderer
    from pptx import Presentation

    out = tmp_path / "output.pptx"
    renderer = ProfileRenderer(TEMPLATE_PATH, SPEC_PATH)
    renderer.render(FIXTURE_PROFILE, out)
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) >= 1


@pytest.mark.skipif(
    not TEMPLATE_PATH.exists(),
    reason="base_profile_template.pptx not yet created (Phase 1 manual step)"
)
def test_no_unreplaced_tokens(tmp_path):
    from core.profile_renderer import ProfileRenderer
    from pptx import Presentation

    out = tmp_path / "output.pptx"
    renderer = ProfileRenderer(TEMPLATE_PATH, SPEC_PATH)
    renderer.render(FIXTURE_PROFILE, out)

    prs = Presentation(str(out))
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = " ".join(
                "".join(r.text for r in p.runs)
                for p in shape.text_frame.paragraphs
            )
            assert "{{" not in full_text, f"Unreplaced token found in shape: {full_text[:100]}"


@pytest.mark.skipif(
    not TEMPLATE_PATH.exists(),
    reason="base_profile_template.pptx not yet created (Phase 1 manual step)"
)
def test_template_not_modified(tmp_path):
    from core.profile_renderer import ProfileRenderer
    import hashlib

    original_hash = hashlib.md5(TEMPLATE_PATH.read_bytes()).hexdigest()
    out = tmp_path / "output.pptx"
    renderer = ProfileRenderer(TEMPLATE_PATH, SPEC_PATH)
    renderer.render(FIXTURE_PROFILE, out)
    after_hash = hashlib.md5(TEMPLATE_PATH.read_bytes()).hexdigest()
    assert original_hash == after_hash, "Template file was modified!"
