import json
import re
from pathlib import Path
from core.profile_schema import ProfileSchema

SYSTEM_PROMPT = """You are a professional profile extraction assistant. Extract structured information from the provided consultant profile source text and return a single valid JSON object.

YOUR CORE JOB:
Read and deeply understand the full source text. Your primary goal is to fill the KEY PROJECTS section as richly as possible using every relevant fact in the source. You are an intelligent synthesizer, not a copy-paste tool. The only hard rule: do NOT invent facts not present in the source.

GENERAL RULES:
- Do NOT use first-person language (no "I", "my", "we").
- Do NOT use: passionate, enthusiast, thrive, committed, driven, dedicated, love, enjoy.
- Do NOT add marketing language or generic filler.
- Profile field: third-person, factual, consultant-style prose. 3-4 sentences. Synthesize the most impressive facts: scale, industries, technologies, leadership scope.
- Technologies list: deduplicated, official product names only.
- Competencies: short labels (2-4 words), not sentences.
- If a field has no data, use empty string or empty list.
- Return ONLY the JSON object. No explanation, no markdown fences, no extra text.

FRAME BUDGET - READ THIS CAREFULLY:
The KEY PROJECTS section uses Arial 8pt font and has space for approximately 36 lines.
Structure overhead (employer headings, project name lines, blank separators) uses ~12 lines.
This leaves ~24 lines of content across all 6 projects = ~4 lines per project.
At Arial 8pt with ~170 chars/line, each project content can hold up to 600-680 characters.
IMPORTANT: Write as much as the source supports. Target 500-650 characters per project when source is rich.
Only write less when the source truly has little information for that project.

CONTENT FORMAT - ADAPTIVE AND MIXED:
Each project has a "content" field. Based on source richness, choose the best format:

FORMAT A - FLOWING PARAGRAPH (preferred when source has rich detail):
  3-4 sentences combining achievements, technologies used, team context, and business impact.
  Written in third-person past tense. No bullet markers. Target 500-650 characters.
  Example: "Led a team of 6 developers in building a master data management platform for a major investment bank, centralising firm-wide reference data distribution across 12 downstream systems. Migrated the application from on-premises infrastructure to GCP using Terraform and Docker, replacing the RabbitMQ messaging layer with Google Pub/Sub. Secured all service-to-service APIs with OAuth 2.0 and JWT tokens, and introduced automated integration testing with JUnit and Mockito reducing regression time by 40%."

FORMAT B - BULLET LINES (use when facts are distinct and don't flow well as prose):
  3-5 lines each starting with \u2022 (bullet) and a space, separated by \n.
  Each line starts with a past-tense verb. Each line 100-140 characters.
  Example: "\u2022 Designed and implemented REST API layer using Spring Boot and Oracle Database with full CI/CD pipeline on Jenkins.\n\u2022 Led architectural reviews and code quality sessions for a team of 4, reducing critical defects by 30%.\n\u2022 Integrated SonarQube static analysis and automated unit test coverage reporting."

FORMAT C - MIXED (use when you have a strong narrative AND additional distinct facts worth preserving):
  Start with 2 sentences of paragraph prose, then add 2-3 bullet lines separated by \n.
  Example: "Developed a workflow management tool for task tracking and process automation across 3 business units, built on Spring Boot with REST APIs and Oracle DB.\n\u2022 Deployed to Kubernetes clusters with Jenkins and SonarQube CI/CD pipelines.\n\u2022 Led cross-team coordination between backend, frontend, and QA.\n\u2022 Delivered 2 major releases within budget, reducing manual processing time by 25%."

Choose the format that best uses available source material. Maximize content within the 650 character limit.
Do NOT write short summaries when more detail exists in the source.

EXPERIENCE PRIORITY STRATEGY:
1. MOST RECENT EMPLOYER: up to 4 named projects, each with full content (aim for 500-650 chars each).
2. SECOND MOST RECENT EMPLOYER: up to 2 projects, each with content (aim for 400-550 chars each).
3. ALL REMAINING EMPLOYERS: one aggregated entry:
   - employer: "Previous Experience"
   - role: one-line summary (e.g. "Senior Developer, Tech Lead - multiple clients")
   - date_range: earliest to latest year
   - employer_bullets: 2-3 synthesized lines (use adaptive format)
   - projects: empty list

PROJECT NAMING:
- Use explicit project names from source when available.
- If no explicit names, group activities into 2-4 meaningful named themes.
- Every non-aggregated employer must have at least 2 projects.
- Never leave content empty if achievements or activities are present in the source.

STRICT LIMITS:
  * competencies: max 8
  * technologies: max 15, deduplicated
  * methodologies: max 5
  * education: max 2
  * certifications: max 3
  * experience entries: exactly 3
  * employer_bullets: max 2 for top 2 entries; max 3 for aggregated
  * projects entry 1: max 4
  * projects entry 2: max 2
  * content per project: max 650 characters
  * role_subtitle: max 5 words

REQUIRED JSON SCHEMA:
{
  "name": "string",
  "role_title": "string",
  "role_subtitle": "string",
  "profile": "string",
  "competencies": ["string"],
  "education": ["string"],
  "certifications": ["string"],
  "methodologies": ["string"],
  "technologies": ["string"],
  "experience": [
    {
      "employer": "string",
      "role": "string",
      "date_range": "string",
      "employer_bullets": ["string"],
      "projects": [
        {
          "project_name": "string",
          "date_range": "string",
          "content": "string"
        }
      ]
    }
  ]
}"""


def _cell_text(cell) -> str:
    lines = [p.text.strip() for p in cell.paragraphs if p.text.strip()]
    return "\n".join(lines)


def _extract_docx_text(file_path: Path) -> str:
    from docx import Document
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    doc = Document(str(file_path))
    output_lines = []
    emitted_tc_ids: set = set()

    def emit_cell(cell):
        tc_id = id(cell._tc)
        if tc_id in emitted_tc_ids:
            return None
        emitted_tc_ids.add(tc_id)
        return _cell_text(cell) or None

    body = doc.element.body
    for child in body.iterchildren():
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'p':
            para = Paragraph(child, doc)
            text = para.text.strip()
            if text:
                output_lines.append(text)
        elif tag == 'tbl':
            tbl = Table(child, doc)
            row_texts = []
            for row in tbl.rows:
                row_parts = []
                for cell in row.cells:
                    text = emit_cell(cell)
                    if text:
                        row_parts.append(text)
                if row_parts:
                    row_texts.append(' || '.join(row_parts))
            output_lines.extend(row_texts)

    return "\n".join(output_lines)


def _extract_text_from_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".docx":
        return _extract_docx_text(file_path)
    elif suffix == ".pdf":
        import fitz
        doc = fitz.open(str(file_path))
        pages = [doc.load_page(i).get_text() for i in range(len(doc))]
        return "\n".join(pages)
    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        texts.append("".join(run.text for run in para.runs))
        return "\n".join(texts)
    elif suffix == ".txt":
        text = file_path.read_text(encoding="utf-8")
        return re.sub(r'[ \t]+', ' ', text).strip()
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


class ProfileExtractor:
    def __init__(self, copilot_client):
        self.client = copilot_client

    def extract_from_file(self, file_path: Path) -> ProfileSchema:
        raw_text = _extract_text_from_file(file_path)
        return self.extract(raw_text)

    def extract(self, raw_text: str) -> ProfileSchema:
        response = self.client.complete(SYSTEM_PROMPT, raw_text)
        cleaned = re.sub(r'^```(?:json)?\s*', '', response.strip())
        cleaned = re.sub(r'```\s*$', '', cleaned).strip()
        try:
            data = json.loads(cleaned)
        except json.JSONDecodeError as e:
            raise ValueError(f"AI returned invalid JSON: {e}\n\nRaw response:\n{response[:500]}")
        try:
            return ProfileSchema(**data)
        except Exception as e:
            raise ValueError(f"Schema validation failed: {e}")
