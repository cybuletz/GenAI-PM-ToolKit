import json
import math
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
from core.profile_schema import ProfileSchema

# Frame metrics measured from base_profile_template.pptx
_FRAME_HEIGHT_EMU = 4212238
_FRAME_WIDTH_EMU  = 8884366
_LINE_HEIGHT_EMU  = int(8 * 12700 * 1.15)
_FRAME_CAPACITY   = _FRAME_HEIGHT_EMU // _LINE_HEIGHT_EMU  # 36 lines
_CHARS_PER_LINE   = int((_FRAME_WIDTH_EMU / 914400) * 17)  # ~165 chars/line

# Global font applied to all generated/replaced text
_FONT             = "Arial"
_CONTENT_SZ       = 800   # 8pt in half-points

# Tokens that should have Arial applied after replacement
_FONT_PATCHED_TOKENS = {
    "{{PROFILE}}",
    "{{EDUCATION}}",
    "{{METHODOLOGIES}}",
    "{{TECHNOLOGIES}}",
    "{{ROLE_TITLE}}",
    "{{NAME}}",
    "{{ROLE_SUBTITLE}}",
    "{{COMP_1}}", "{{COMP_2}}", "{{COMP_3}}", "{{COMP_4}}",
    "{{COMP_5}}", "{{COMP_6}}", "{{COMP_7}}", "{{COMP_8}}",
}

_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class ProfileRenderer:
    def __init__(self, template_path: Path, spec_path: Path):
        self.template_path = template_path
        self.spec = json.loads(spec_path.read_text(encoding="utf-8"))

    def render(self, profile: ProfileSchema, output_path: Path):
        prs = Presentation(str(self.template_path))
        slide_index = self.spec.get("slide_index", 0)
        slide = prs.slides[slide_index]
        text_replacements = self._build_text_replacements(profile)
        projects_lines = self._build_projects_block(profile.experience)
        # Track which shapes had tokens replaced so we can patch fonts
        patched_shapes = set()
        for shape in slide.shapes:
            self._process_shape(shape, text_replacements, projects_lines, patched_shapes)
        # Apply Arial to all runs in shapes that had token replacements
        for shape in slide.shapes:
            if id(shape) in patched_shapes:
                self._patch_font_in_shape(shape)
        prs.save(str(output_path))

    def _process_shape(self, shape, text_replacements, projects_lines, patched_shapes):
        if shape.shape_type == 6:
            try:
                for child in shape.shapes:
                    self._process_shape(child, text_replacements, projects_lines, patched_shapes)
            except Exception:
                pass
            return
        if shape.has_text_frame:
            replaced = self._replace_in_shape(shape, text_replacements, projects_lines)
            if replaced:
                patched_shapes.add(id(shape))

    def _build_text_replacements(self, profile: ProfileSchema) -> dict:
        r = {}
        r["{{ROLE_TITLE}}"] = profile.role_title
        r["{{NAME}}"] = profile.name
        r["{{ROLE_SUBTITLE}}"] = profile.role_subtitle
        r["{{PROFILE}}"] = profile.profile
        for i in range(8):
            r[f"{{{{COMP_{i+1}}}}}"] = profile.competencies[i] if i < len(profile.competencies) else ""
        r["{{EDUCATION}}"] = " | ".join(profile.education)
        r["{{METHODOLOGIES}}"] = ", ".join(profile.methodologies)
        r["{{TECHNOLOGIES}}"] = ", ".join(profile.technologies)
        return r

    def _patch_font_in_shape(self, shape):
        """Set Arial on every run in the shape's text frame."""
        if not shape.has_text_frame:
            return
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = etree.SubElement(run._r, qn('a:rPr'))
                    run._r.insert(0, rPr)  # rPr must be first child
                # Remove existing latin/ea/cs font elements
                for tag in ['a:latin', 'a:ea', 'a:cs']:
                    existing = rPr.find(qn(tag))
                    if existing is not None:
                        rPr.remove(existing)
                # Insert Arial as latin typeface
                latin = etree.SubElement(rPr, qn('a:latin'))
                latin.set('typeface', _FONT)

    def _estimate_lines(self, text: str, bold: bool = False) -> int:
        if not text:
            return 1
        if '\n' in text:
            return max(1, sum(self._estimate_lines(s.strip(), bold) for s in text.split('\n') if s.strip()))
        cpl = int(_CHARS_PER_LINE * 0.82) if bold else _CHARS_PER_LINE
        return max(1, math.ceil(len(text) / cpl))

    def _content_to_line_specs(self, content: str) -> list:
        if not content:
            return []
        return [(seg.strip(), False, _CONTENT_SZ) for seg in content.split('\n') if seg.strip()]

    def _build_projects_block(self, experience: list) -> list:
        lines = []
        used = 0
        capacity = _FRAME_CAPACITY - 1

        for i, exp in enumerate(experience):
            sep_cost = 0 if i == 0 else 1
            heading = f"{exp.employer} \u2013 {exp.role} / {exp.date_range}"
            heading_cost = self._estimate_lines(heading, bold=True)

            if used + sep_cost + heading_cost > capacity:
                break

            if i > 0:
                lines.append(("", False, _CONTENT_SZ))
                used += 1
            lines.append((heading, True, None))
            used += heading_cost

            if exp.projects:
                for proj in exp.projects:
                    content = getattr(proj, 'content', '') or ''
                    content_specs = self._content_to_line_specs(content)
                    content_cost = sum(self._estimate_lines(t) for t, _, _ in content_specs)
                    block_cost = 1 + content_cost
                    if used + block_cost > capacity:
                        break
                    lines.append((proj.project_name, True, _CONTENT_SZ))
                    used += 1
                    lines.extend(content_specs)
                    used += content_cost
            else:
                for b in exp.employer_bullets:
                    b_specs = self._content_to_line_specs(b)
                    b_cost = sum(self._estimate_lines(t) for t, _, _ in b_specs)
                    if used + b_cost > capacity:
                        break
                    lines.extend(b_specs)
                    used += b_cost

        return lines

    def _replace_in_shape(self, shape, text_replacements: dict, projects_lines: list) -> bool:
        """Returns True if any replacement was made."""
        tf = shape.text_frame
        replaced = False
        for para in tf.paragraphs:
            via_xml = "".join(
                (t.text or "") for t in para._p.iter() if t.tag.endswith("}t")
            )
            if "{{" not in via_xml:
                continue
            if "{{KEY_PROJECTS}}" in via_xml:
                self._rebuild_projects_block(para, projects_lines)
                return True
            self._apply_replacements_to_para(para, text_replacements)
            replaced = True
        return replaced

    def _apply_replacements_to_para(self, para, replacements: dict):
        runs = para.runs
        if not runs:
            for t_node in para._p.iter():
                if t_node.tag.endswith("}t") and t_node.text and "{{" in t_node.text:
                    for token, value in replacements.items():
                        t_node.text = t_node.text.replace(token, value)
            return
        combined = "".join(r.text or "" for r in runs)
        if "{{" not in combined:
            return
        for token, value in replacements.items():
            if token not in combined:
                continue
            replaced = False
            for run in runs:
                if token in (run.text or ""):
                    run.text = run.text.replace(token, value)
                    replaced = True
                    break
            if not replaced:
                self._heal_split_token(runs, token, value)
            combined = "".join(r.text or "" for r in runs)

    def _heal_split_token(self, runs, token: str, value: str):
        texts = [r.text or "" for r in runs]
        n = len(texts)
        for start in range(n):
            for end in range(start + 1, n + 1):
                segment = "".join(texts[start:end])
                if token in segment:
                    runs[start].text = segment.replace(token, value)
                    for i in range(start + 1, end):
                        runs[i].text = ""
                    return

    def _rebuild_projects_block(self, anchor_para, line_specs: list):
        tf_elem = anchor_para._p.getparent()
        existing_paras = tf_elem.findall(qn('a:p'))

        base_color = None
        ref_runs = anchor_para._p.findall('.//' + qn('a:r'))
        if ref_runs:
            rPr = ref_runs[0].find(qn('a:rPr'))
            if rPr is not None:
                sf = rPr.find('.//' + qn('a:solidFill'))
                if sf is not None:
                    srgb = sf.find(qn('a:srgbClr'))
                    if srgb is not None:
                        base_color = srgb.get('val')

        for p in existing_paras:
            tf_elem.remove(p)

        def _make_p(text, bold, sz):
            color_xml = f'<a:solidFill><a:srgbClr val="{base_color}"/></a:solidFill>' if base_color else ""
            sz_xml = f' sz="{sz}"' if sz is not None else ''
            b_xml = ' b="1"' if bold else ' b="0"'
            safe = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            return parse_xml(
                f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:pPr><a:lnSpc><a:spcPct val="100000"/></a:lnSpc></a:pPr>'
                f'<a:r><a:rPr lang="en-US"{sz_xml}{b_xml} dirty="0">'
                f'{color_xml}<a:latin typeface="Arial"/>'
                f'</a:rPr>'
                f'<a:t>{safe}</a:t></a:r></a:p>'
            )

        prev_p = None
        for (text, bold, sz) in line_specs:
            new_p = _make_p(text, bold, sz)
            if prev_p is None:
                tf_elem.append(new_p)
            else:
                prev_p.addnext(new_p)
            prev_p = new_p
