"""Run this script once on your prepared template to verify all tokens
are readable and show exactly how they are stored in XML runs.

Usage:
  cd pm-toolkit
  python tools/diagnose_template.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

TEMPLATE = Path(__file__).parent.parent / "templates" / "base_profile_template.pptx"

EXPECTED_TOKENS = [
    "{{ROLE_TITLE}}", "{{NAME}}", "{{ROLE_SUBTITLE}}",
    "{{COMP_1}}", "{{COMP_2}}", "{{COMP_3}}", "{{COMP_4}}",
    "{{COMP_5}}", "{{COMP_6}}", "{{COMP_7}}", "{{COMP_8}}",
    "{{PROFILE}}", "{{KEY_PROJECTS}}",
    "{{EDUCATION}}", "{{METHODOLOGIES}}", "{{TECHNOLOGIES}}",
]


def all_text_from_element(elem) -> str:
    """Collect ALL text from every <a:t> node anywhere in elem's subtree."""
    return "".join((t.text or "") for t in elem.iter() if t.tag.endswith("}t"))


def scan_shape(shape, found_tokens, indent=""):
    """Recursively scan a shape (handles groups)."""
    # Groups: recurse into children
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP = 6
        try:
            for child in shape.shapes:
                scan_shape(child, found_tokens, indent + "  ")
        except Exception:
            pass
        return

    if not shape.has_text_frame:
        # Still dump raw XML text to catch any edge cases
        raw = all_text_from_element(shape._element)
        if "{{" in raw:
            print(f"{indent}Shape {shape.shape_id} ({shape.name!r}) [NO text_frame]")
            print(f"{indent}  raw XML text: {raw!r}")
            for token in EXPECTED_TOKENS:
                if token in raw:
                    found_tokens.add(token)
        return

    tf = shape.text_frame
    for para in tf.paragraphs:
        via_runs = "".join(r.text for r in para.runs)
        via_xml = "".join((t.text or "") for t in para._p.iter() if t.tag.endswith("}t"))

        if "{{" not in via_xml and "{{" not in via_runs:
            continue

        print(f"{indent}Shape {shape.shape_id} ({shape.name!r})")
        print(f"{indent}  via .runs : {via_runs!r}")
        print(f"{indent}  via XML   : {via_xml!r}")
        print(f"{indent}  run count : {len(para.runs)}")
        for i, run in enumerate(para.runs):
            print(f"{indent}  run[{i}]    : {run.text!r}")

        for token in EXPECTED_TOKENS:
            if token in via_xml or token in via_runs:
                found_tokens.add(token)


prs = Presentation(str(TEMPLATE))
slide = prs.slides[0]

found_tokens = set()

print("=" * 70)
print("TEMPLATE TOKEN DIAGNOSTIC")
print("=" * 70)

for shape in slide.shapes:
    scan_shape(shape, found_tokens)

print("\n" + "=" * 70)
print("TOKENS FOUND IN TEMPLATE:")
for t in EXPECTED_TOKENS:
    status = "OK" if t in found_tokens else "MISSING"
    print(f"  [{status:7s}] {t}")

missing = [t for t in EXPECTED_TOKENS if t not in found_tokens]
if missing:
    print(f"\n  WARNING: {len(missing)} token(s) not found.")
    print("  These tokens must be added to the template manually in PowerPoint.")
    print("  Open the template, find the corresponding shape, and type the token exactly.")
else:
    print("\n  All tokens present. Template is ready.")

# --- BONUS: dump ALL shapes with any text, token or not ---
print("\n" + "=" * 70)
print("ALL SHAPES WITH TEXT (for reference):")
for shape in slide.shapes:
    if shape.shape_type == 6:
        print(f"  GROUP {shape.shape_id} ({shape.name!r})")
        try:
            for child in shape.shapes:
                raw = all_text_from_element(child._element)
                if raw.strip():
                    print(f"    child {child.shape_id} ({child.name!r}): {raw.strip()[:80]!r}")
        except Exception:
            pass
        continue
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    if text:
        print(f"  Shape {shape.shape_id} ({shape.name!r}): {text[:80]!r}")
